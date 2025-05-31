# from speakListen import hear
import sys
import os
# Get the absolute path to the project root directory
project_root = os.path.abspath(os.path.join(os.path.dirname(__file__), '..'))
# Add the project root to the Python path
sys.path.insert(0, project_root)
# from ENGINE.STT.fast_stt import speech_to_text
from ENGINE.STT.apple_stt import speech_to_text # <-- Use Apple STT
# from ENGINE.TTS.fast_TTS_DF import speak
# from ENGINE.TTS.TTS_DF import speak
from ENGINE.TTS.eSpeakNG_fast50ms import speak
import docx
import fitz
import time
from rich.console import Console # pip3 install Rich
from rich.table import Table
from colorama import Fore
from pptx import Presentation
import subprocess
from BRAIN.ai_chat_res.functions_call import get_input_from_popup

def get_file_path(location):
    """Search for the file in the current directory, as an absolute path, or recursively in the system.

    Args:
        location (str): File name or full path.

    Returns:
        str: Full file path if found, otherwise None.
    """
    # Check if the input is a valid file path
    if os.path.isfile(location):
        return os.path.abspath(location)
    
    # Check in the current working directory
    current_dir = os.getcwd()
    potential_path = os.path.join(current_dir, location)
    if os.path.isfile(potential_path):
        return os.path.abspath(potential_path)
    
    # If not found, search recursively in the user's home directory
    print(Fore.YELLOW + "Searching for the file across the system. This may take some time...")
    home_dir = os.path.expanduser("~")  # User's home directory
    for root, dirs, files in os.walk(home_dir):
        if location in files:
            return os.path.join(root, location)
    
    # If the file is not found
    return None

def ms_word(location: str):
    """Print and speak out an MS Word docx file as specified in the path or name."""
    try:
        # speak("Enter the document's name or full location - ")
        # location = input("Enter the document's name or full location - ")
        
        file_loc = get_file_path(location)
        if not file_loc:
            print(Fore.YELLOW + "I couldn't locate the file! Please check the name or path.")
            speak("I couldn't locate the file! Please check the name or path.")
            return "File not found. Please check the name or path."
       
        doc = docx.Document(file_loc)
        fullText = []
        for para in doc.paragraphs:
            fullText.append(para.text)
        doc_file = '\n'.join(fullText)
        print(doc_file)
        speak(doc_file)
        return f"Successfully read Word document: {os.path.basename(file_loc)}"
    except Exception as exp:
        print(f"ERROR - {exp}")
        print(Fore.YELLOW + "I couldn't process the file. Please check the format or content.")
        return f"Error processing Word document: {exp}"
def pdf_read(location: str):
    """Print and speak out the PDF on the specified path or name."""
    try:
        # speak("Enter the document's name or full location - ")
        # location = input("Enter the document's name or full location - ")
        
        # Handle file path for both macOS and Windows
        file_loc = get_file_path(location)
        if not file_loc:
            print(Fore.YELLOW + "I couldn't locate the file! Please check the name or path.")
            speak("I couldn't locate the file! Please check the name or path.")
            return "PDF file not found. Please check the name or path."
        
        pdf = fitz.open(file_loc)
        details = pdf.metadata  # Metadata includes Author name and Title of book/document.
        total_pages = pdf.page_count  # Total number of pages

        # Extract metadata
        author = details.get("author", "Unknown")
        title = details.get("title", "Unknown")
        book_details(author, title, total_pages)
        speak(f"Title: {title}")
        speak(f"Author: {author}")
        speak(f"Total Pages: {total_pages}")
        
        # Handle the Index
        toc = pdf.get_toc()
        print("Say 1 or \"ONLY PRINT INDEX\" - if you want me to print the book's index.\n"
              "Say 2 if you want me to print and make me speak out the book's index.\n"
              "Say any key if you don't want to print the index.")
        speak("Say 1 or only print index if you want me to print the book's index.\n"
              "Say 2 if you want me to print and make me speak out the book's index.\n"
              "Say any key if you don't want to print the index.")
        q = speech_to_text().lower()

        if "only print" in q or "1" in q or "One" in q or "one" in q:
            print_index(toc)
        elif "speak" in q or "2" in q or "Two" in q or "two" in q:
            print_n_speak_index(toc)
        elif q == "none":
            print("I couldn't understand what you just said!")
            speak("I couldn't understand what you just said!")
        else:
            pass

        # Allow the user to choose reading options
        print("____________________________________________________________________________________________________________")
        print("1. Print/speak a single page\n2. Print/speak a range of pages\n3. Print/speak a Lesson\n4. Read/speak a whole book")
        speak("1. Print/speak a single page\n2. Print/speak a range of pages\n3. Print/speak a Lesson\n4. Read/speak a whole book")
        q = speech_to_text().lower()

        if "single" in q or "one" in q or "1" in q or "One" in q:
            try:
                speak("Enter the page number")
                pgno = int(get_input_from_popup("Page Number : "))
                page = pdf.load_page(pgno - 1)
                text = page.get_text('text')
                print(text.replace('\t', ' '))
                speak(text.replace('\t', ' '))
            except Exception:
                print("Sorry, I couldn't recognize what you entered. Please re-enter the Page Number.")
                speak("Sorry, I couldn't recognize what you entered. Please re-enter the Page Number.")

        elif "range" in q or "multiple" in q or "2" in q or "Two" in q or "two" in q:
            try:
                speak("Enter the starting and ending page numbers")
                start_pg_no = int(get_input_from_popup("Starting Page Number : "))
                end_pg_no = int(get_input_from_popup("End Page Number - "))
                for i in range(start_pg_no - 1, end_pg_no):
                    page = pdf.load_page(i)
                    text = page.get_text('text')
                    print(text.replace('\t', ' '))
                    speak(text.replace('\t', ' '))
            except Exception:
                print("Sorry, I couldn't recognize what you entered. Please re-enter the Page Numbers.")
                speak("Sorry, I couldn't recognize what you entered. Please re-enter the Page Numbers.")

        elif "lesson" in q or "3" in q:
            try:
                key = get_input_from_popup("Lesson name - ")
                start_pg_no, end_pg_no = search_in_toc(toc, key, total_pages)
                if start_pg_no is not None and end_pg_no is not None:
                    for i in range(start_pg_no - 1, end_pg_no):
                        page = pdf.load_page(i)
                        text = page.get_text('text')
                        print(text.replace('\t', ' '))
                        speak(text.replace('\t', ' '))
                else:
                    print("Sorry, I cannot find the particular lesson.")
                    speak("Sorry, I cannot find the particular lesson.")
            except Exception:
                print("Try Again! Lesson could not be found.")
                speak("Try Again! Lesson could not be found.")

        elif "whole" in q or "complete" in q or "4" in q:
            for i in range(total_pages):
                page = pdf.load_page(i)
                text = page.get_text('text')
                print(text.replace('\t', ' '))
                speak(text.replace('\t', ' '))

        elif q == "none":
            print("I couldn't understand what you just said!")
            speak("I couldn't understand what you just said!")
        else:
            print("You didn't say a valid command!")
            time.sleep(5)
        return f"Finished interacting with PDF: {os.path.basename(file_loc)}"

    except Exception as e:
        print(f"ERROR - {e}")
        print(Fore.YELLOW + "I couldn't process the file. Please check the format or content.")
        return f"Error during PDF processing: {e}"
    finally:
        if 'pdf' in locals() and pdf and not pdf.is_closed: # Ensure pdf exists and is open
            pdf.close()

def presentation_read(location: str):
    """Read and speak out the content of a PowerPoint (.pptx) or Keynote (.key) file."""
    try:
        # speak("Enter the presentation's name or full location - ")
        # location = input("Enter the presentation's name or full location - ")

        # Search for the file
        file_loc = get_file_path(location)
        if not file_loc:
            print(Fore.YELLOW + "I couldn't locate the file! Please check the name or path.")
            speak("I couldn't locate the file! Please check the name or path.")
            return "Presentation file not found. Please check the name or path."

        # Determine file type
        if file_loc.endswith(".pptx"):
            return pptx_read(file_loc)
        elif file_loc.endswith(".key"):
            if sys.platform == "darwin":  # macOS
                return key_read(file_loc) # Return the result from key_read
            else:
                msg = "Keynote files are not supported on this OS. Please provide a .pptx file."
                print(Fore.YELLOW + msg)
                speak(msg)
                return msg
        else:
            msg = "Unsupported presentation file format. Please provide a .pptx or .key file."
            print(Fore.YELLOW + msg)
            speak(msg)
            return msg
    except Exception as exp:
        print(f"ERROR - {exp}")
        print(Fore.YELLOW + "I couldn't process the presentation file. Please check the format or content.")
        speak("I couldn't process the presentation file. Please check the format or content.")
        return f"Error processing presentation: {exp}"

def pptx_read(file_loc):
    """Read and speak out the content of a PowerPoint (.pptx) file."""
    try:
        presentation = Presentation(file_loc)
        full_text = []

        # Extract text from each slide
        for slide_number, slide in enumerate(presentation.slides, start=1):
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        slide_text.append(paragraph.text)
            slide_content = "\n".join(slide_text)
            full_text.append(f"Slide {slide_number}:\n{slide_content}")

        # Combine all slides' text
        ppt_content = "\n\n".join(full_text)
        print(ppt_content)
        speak(ppt_content)
        return f"Successfully read PowerPoint presentation: {os.path.basename(file_loc)}"
    except Exception as exp:
        print(f"ERROR - {exp}")
        print(Fore.YELLOW + "I couldn't process the PowerPoint file. Please check the format or content.")
        return f"Error processing PowerPoint file: {exp}"

def key_read(file_loc):
    """Read and speak out the content of a Keynote (.key) file using AppleScript (macOS only)."""
    try:
        # Convert the input file location to an absolute POSIX path
        posix_path = os.path.abspath(file_loc)

        # AppleScript to open the Keynote file and extract text
        script = f'''
        tell application "Keynote"
            open POSIX file "{posix_path}"
            set theDoc to front document
            set slideTexts to ""
            repeat with eachSlide in slides of theDoc
                repeat with eachShape in shapes of eachSlide
                    if class of eachShape is text then
                        set slideTexts to slideTexts & (content of text of eachShape) & linefeed
                    end if
                end repeat
                set slideTexts to slideTexts & linefeed
            end repeat
            close theDoc
            return slideTexts
        end tell
        '''

        # Execute the AppleScript
        result = subprocess.run(
            ["osascript", "-e", script], 
            capture_output=True, 
            text=True
        )

        if result.returncode == 0:
            content = result.stdout.strip()
            if content:
                print(content)
                speak(content)
                return f"Successfully read Keynote presentation: {os.path.basename(file_loc)}"
            else:
                msg = "The Keynote file contains no text."
                print(Fore.YELLOW + msg)
                speak(msg)
                return msg
        else:
            print(Fore.RED + "Error executing AppleScript:", result.stderr)
            msg = "I couldn't process the Keynote file via AppleScript. Please check the format or content."
            print(Fore.YELLOW + msg)
            speak(msg)
            return msg

    except Exception as exp:
        print(f"ERROR - {exp}")
        return f"Error processing Keynote file: {exp}"

def doubleslash(text):
    """Replaces / with // 

    Args:
        text (str): location

    Returns:
        str: formatted location
    """
    return text.replace('\\' , '\\\\')

def print_index(toc):
    """Prints out the index in proper format with title name and page number

    Args:
        toc (nested list): toc[1] - Topic name
                           toc[2] - Page number
    """
    dash = "-"*(100 - 7)
    space = " "*47
    print(f"{space}INDEX")
    print(f"\n\nName : {dash} PageNo.\n\n\n")
    for topic in toc:
        eq_dash = "-"*(100 - len(topic[1]))
        print(f"{topic[1]} {eq_dash} {topic[2]}")
        
def print_n_speak_index(toc):
    """Along with printing, it speaks out the index too.

    Args:
        toc (nested list): toc[1] - Topic name
                           toc[2] - Page number
    """
    dash = "-"*(100 - 7)
    space = " "*47
    print(f"{space}INDEX")
    print(f"\n\nName : {dash} PageNo.\n\n\n\n")
    for topic in toc:
        eq_dash = "-"*(100 - len(topic[1]))
        print(f"{topic[1]} {eq_dash} {topic[2]}")
        speak(f"{topic[1]}  {topic[2]}")

def search_in_toc(toc, key, totalpg):
    """Searches a particular lesson name provided as a parameter in toc and returns its starting and ending page numbers.

    Args:
        toc (nested list): toc[1] - Topic name
                           toc[2] - Page number
        key (str): the key to be found
        totalpg (int): total pages in book/document

    Returns:
        int: staring and ending page numbers of lesson found.
        If not found then return None
    """
    for i in range(len(toc) - 1):
        topic = toc[i]
        if i != len(toc) - 2:
            if topic[1] == key:
                nexttopic = toc[i + 1]
                return (topic[2], nexttopic[2])
            elif topic[1].lower() == key:
                nexttopic = toc[i + 1]
                return (topic[2], nexttopic[2])
        else:
            if topic[1] == key:
                return (topic[2], totalpg)
            elif topic[1].lower() == key:
               
                return (topic[2], totalpg)
    return None,None

def book_details(author, title, total_pages):
    """Creates a table of book details like author name, title, and total pages.

    Args:
        author (str): Name of author
        title (str): title of the book
        total_pages (int): total pages in the book
    """
    table = Table(title="\nBook Details :- ", show_lines = True) 

    table.add_column("Sr. No.", style="magenta", no_wrap=True)
    table.add_column("Property", style="cyan")
    table.add_column("Value", justify="left", style="green")

    table.add_row("1", "Title", f"{title}")
    table.add_row("2", "Author", f"{author}")
    table.add_row("3", "Pages", f"{total_pages}")

    console = Console()
    console.print(table)
   
#ms_word()
# pdf_read()
#book_details("abc", "abcde", 12)
# presentation_read()

